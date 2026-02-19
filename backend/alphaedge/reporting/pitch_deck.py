"""Generate professional PPTX pitch decks from analysis results."""

from __future__ import annotations

import logging
import os
import shutil
import tempfile
from io import BytesIO
from typing import Any

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x14, 0x2B)
DARK_PANEL = RGBColor(0x10, 0x1D, 0x3A)
CARD_BG = RGBColor(0x14, 0x22, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xE2, 0xE8, 0xF0)
LIGHT_GRAY = RGBColor(0xA0, 0xAE, 0xC4)
MID_GRAY = RGBColor(0x64, 0x74, 0x8B)
DIM_GRAY = RGBColor(0x47, 0x55, 0x69)
ACCENT_BLUE = RGBColor(0x63, 0x66, 0xF1)   # indigo-500
ACCENT_BLUE_DIM = RGBColor(0x45, 0x48, 0xB0)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_GREEN_DIM = RGBColor(0x0D, 0x8A, 0x62)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT_GOLD = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_CYAN = RGBColor(0x06, 0xB6, 0xD4)
BORDER_COLOR = RGBColor(0x1E, 0x29, 0x3B)
BORDER_LIGHT = RGBColor(0x33, 0x41, 0x55)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# matplotlib dark theme
_MPL_RC = {
    "figure.facecolor": "#0B142B",
    "axes.facecolor": "#101D3A",
    "text.color": "#FFFFFF",
    "axes.labelcolor": "#A0AEC4",
    "xtick.color": "#A0AEC4",
    "ytick.color": "#A0AEC4",
    "axes.edgecolor": "#1E293B",
    "grid.color": "#1E293B",
    "grid.alpha": 0.5,
    "font.family": "sans-serif",
    "font.size": 11,
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _sg(data: dict | None, *keys: str, default: Any = None) -> Any:
    """Safely traverse nested dicts."""
    cur = data
    for k in keys:
        if cur is None or not isinstance(cur, dict):
            return default
        cur = cur.get(k)
    return cur if cur is not None else default


def _fmt_large(n: Any) -> str:
    if n is None:
        return "\u2014"
    try:
        n = float(n)
    except (TypeError, ValueError):
        return "\u2014"
    if abs(n) >= 1e12:
        return f"${n / 1e12:.2f}T"
    if abs(n) >= 1e9:
        return f"${n / 1e9:.2f}B"
    if abs(n) >= 1e6:
        return f"${n / 1e6:.1f}M"
    if abs(n) >= 1e3:
        return f"${n / 1e3:.1f}K"
    return f"${n:,.0f}"


def _fmt_pct(n: Any, decimals: int = 1) -> str:
    if n is None:
        return "\u2014"
    try:
        v = float(n)
        # If already in % form (>1 or <-1 and not a ratio), display directly
        if abs(v) <= 1.0:
            return f"{v * 100:.{decimals}f}%"
        return f"{v:.{decimals}f}%"
    except (TypeError, ValueError):
        return "\u2014"


def _fmt_ratio(n: Any, decimals: int = 2) -> str:
    if n is None:
        return "\u2014"
    try:
        return f"{float(n):.{decimals}f}"
    except (TypeError, ValueError):
        return "\u2014"


def _fmt_price(n: Any) -> str:
    if n is None:
        return "\u2014"
    try:
        return f"${float(n):,.2f}"
    except (TypeError, ValueError):
        return "\u2014"


def _pct_raw(n: Any) -> float | None:
    """Return a percentage as a float (0-100 scale), or None."""
    if n is None:
        return None
    try:
        v = float(n)
        return v * 100 if abs(v) <= 1.0 else v
    except (TypeError, ValueError):
        return None


# ---------------------------------------------------------------------------
# PitchDeckGenerator
# ---------------------------------------------------------------------------

class PitchDeckGenerator:
    """Generate a PPTX pitch deck from AlphaEdge analysis results."""

    def __init__(self) -> None:
        self._prs: Presentation | None = None
        self._charts_dir: str = ""

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def generate(self, analysis: dict) -> BytesIO:
        """Build pitch deck and return as in-memory bytes."""
        self._charts_dir = tempfile.mkdtemp(prefix="alphaedge_deck_")
        try:
            self._prs = Presentation()
            self._prs.slide_width = SLIDE_W
            self._prs.slide_height = SLIDE_H

            self._add_recommendation_slide(analysis)
            self._add_business_overview_slide(analysis)
            self._add_financial_summary_slide(analysis)
            self._add_industry_slide(analysis)
            self._add_why_mispriced_slide(analysis)
            self._add_thesis_slides(analysis)
            self._add_downside_slide(analysis)
            self._add_valuation_slide(analysis)
            self._add_risks_slide(analysis)
            self._add_conclusion_slide(analysis)

            buf = BytesIO()
            self._prs.save(buf)
            buf.seek(0)
            return buf
        finally:
            shutil.rmtree(self._charts_dir, ignore_errors=True)

    # ------------------------------------------------------------------
    # Layout helpers
    # ------------------------------------------------------------------

    def _blank_slide(self):
        layout = self._prs.slide_layouts[6]  # blank
        slide = self._prs.slides.add_slide(layout)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = NAVY
        return slide

    def _title_bar(self, slide, text: str, subtitle: str = ""):
        """Accent title bar at the top with optional subtitle."""
        # accent line
        ln = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.45), Inches(0.08), Inches(0.35),
        )
        ln.fill.solid()
        ln.fill.fore_color.rgb = ACCENT_BLUE
        ln.line.fill.background()
        # title text
        self._text_box(slide, Inches(0.9), Inches(0.35), Inches(10), Inches(0.55),
                       text, size=24, bold=True, color=WHITE)
        if subtitle:
            self._text_box(slide, Inches(0.9), Inches(0.78), Inches(10), Inches(0.35),
                           subtitle, size=11, color=MID_GRAY)

    def _text_box(self, slide, left, top, width, height, text: str, *,
                  size: int = 14, bold: bool = False, color=LIGHT_GRAY,
                  align=PP_ALIGN.LEFT, font_name: str = "Calibri",
                  line_spacing: float | None = None):
        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
        return txbox

    def _rich_text_box(self, slide, left, top, width, height, paragraphs: list[dict]):
        """Add a text box with multiple styled paragraphs.

        Each paragraph dict: {text, size, bold, color, spacing_after, bullet}
        """
        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True

        for i, para in enumerate(paragraphs):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = para.get("text", "")
            p.font.size = Pt(para.get("size", 12))
            p.font.bold = para.get("bold", False)
            p.font.color.rgb = para.get("color", LIGHT_GRAY)
            p.font.name = para.get("font_name", "Calibri")
            p.alignment = para.get("align", PP_ALIGN.LEFT)
            if para.get("spacing_after"):
                p.space_after = Pt(para["spacing_after"])
            if para.get("spacing_before"):
                p.space_before = Pt(para["spacing_before"])
        return txbox

    def _panel(self, slide, left, top, width, height, *, fill=DARK_PANEL,
               border=BORDER_COLOR, radius=True):
        """Add a rounded rectangle panel/card background."""
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        box = slide.shapes.add_shape(shape_type, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = border
        box.line.width = Pt(0.75)
        return box

    def _section_header(self, slide, left, top, text: str, width=Inches(12)):
        """Small uppercase section header with accent underline."""
        self._text_box(slide, left, top, width, Inches(0.25),
                       text.upper(), size=9, bold=True, color=MID_GRAY)
        # thin underline
        ln = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top + Inches(0.28), Inches(1.2), Pt(1.5),
        )
        ln.fill.solid()
        ln.fill.fore_color.rgb = ACCENT_BLUE_DIM
        ln.line.fill.background()

    def _metric_card(self, slide, left, top, label: str, value: str,
                     width=Inches(2.4), height=Inches(0.9), *,
                     value_color=WHITE, label_color=MID_GRAY,
                     subtitle: str = "", subtitle_color=DIM_GRAY):
        """Render a metric card with label, value, and optional subtitle."""
        self._panel(slide, left, top, width, height)
        # label
        self._text_box(slide, left + Inches(0.15), top + Inches(0.08),
                       width - Inches(0.3), Inches(0.22),
                       label.upper(), size=8, bold=True, color=label_color)
        # value
        self._text_box(slide, left + Inches(0.15), top + Inches(0.3),
                       width - Inches(0.3), Inches(0.35),
                       value, size=18, bold=True, color=value_color)
        # subtitle
        if subtitle:
            self._text_box(slide, left + Inches(0.15), top + Inches(0.65),
                           width - Inches(0.3), Inches(0.2),
                           subtitle, size=8, color=subtitle_color)

    def _large_metric_card(self, slide, left, top, label: str, value: str,
                           width=Inches(3.5), height=Inches(1.4), *,
                           value_color=WHITE, description: str = ""):
        """Larger metric card for hero numbers."""
        self._panel(slide, left, top, width, height)
        self._text_box(slide, left + Inches(0.2), top + Inches(0.1),
                       width - Inches(0.4), Inches(0.25),
                       label.upper(), size=9, bold=True, color=MID_GRAY)
        self._text_box(slide, left + Inches(0.2), top + Inches(0.4),
                       width - Inches(0.4), Inches(0.55),
                       value, size=28, bold=True, color=value_color)
        if description:
            self._text_box(slide, left + Inches(0.2), top + Inches(1.0),
                           width - Inches(0.4), Inches(0.3),
                           description, size=9, color=LIGHT_GRAY)

    def _bullet_list(self, slide, left, top, width, height, items: list[str], *,
                     size: int = 12, color=LIGHT_GRAY, spacing: int = 8,
                     bullet: str = "\u25B8"):
        """Render a bullet list with consistent styling."""
        paras = []
        for item in items:
            paras.append({
                "text": f"{bullet}  {item}",
                "size": size,
                "color": color,
                "spacing_after": spacing,
            })
        if paras:
            self._rich_text_box(slide, left, top, width, height, paras)

    def _footer(self, slide, ticker: str = ""):
        """Footer with branding and disclaimer."""
        self._text_box(
            slide, Inches(0.6), Inches(6.95), Inches(6), Inches(0.35),
            "AlphaEdge Analysis Platform  \u2022  For informational purposes only",
            size=7, color=MID_GRAY,
        )
        if ticker:
            self._text_box(
                slide, Inches(10), Inches(6.95), Inches(2.7), Inches(0.35),
                ticker, size=7, color=MID_GRAY, align=PP_ALIGN.RIGHT,
            )

    def _add_picture(self, slide, path: str, left, top, width=None, height=None):
        kwargs = {"left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        slide.shapes.add_picture(path, **kwargs)

    def _divider(self, slide, top):
        ln = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(12.13), Pt(1),
        )
        ln.fill.solid()
        ln.fill.fore_color.rgb = BORDER_COLOR
        ln.line.fill.background()

    def _side_accent(self, slide, top, height, color=ACCENT_BLUE):
        """Vertical accent bar on the left side of a content section."""
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(0.06), height,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

    # ------------------------------------------------------------------
    # SLIDE 1: Recommendation
    # ------------------------------------------------------------------

    def _add_recommendation_slide(self, analysis: dict):
        slide = self._blank_slide()
        snap = analysis.get("snapshot") or {}
        fund = analysis.get("fundamentals") or {}
        verdict = _sg(fund, "verdict") or {}
        thesis = _sg(fund, "investment_thesis") or {}
        combined = fund.get("combined_range") or {}
        forecast = analysis.get("forecast") or {}

        name = snap.get("name", analysis.get("ticker", ""))
        ticker = snap.get("ticker", analysis.get("ticker", ""))
        price = snap.get("price")
        label = verdict.get("label", "")
        upside = verdict.get("upside_pct")
        fair_mid = verdict.get("fair_value_mid")
        fair_low = combined.get("low")
        fair_high = combined.get("high")

        is_long = label in ("undervalued", "fairly_valued")
        rec_text = "LONG" if is_long else "SHORT"
        rec_color = ACCENT_GREEN if is_long else ACCENT_RED

        # Top-left: Recommendation badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(0.8), Inches(2.4), Inches(0.75),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = rec_color
        badge.line.fill.background()
        tf = badge.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = rec_text
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(4)

        # Top-right: Sector tag
        sector = snap.get("sector", "")
        industry = snap.get("industry", "")
        if sector or industry:
            tag_text = f"{sector}  \u2022  {industry}" if sector and industry else (sector or industry)
            self._text_box(slide, Inches(7), Inches(0.9), Inches(5.8), Inches(0.35),
                           tag_text, size=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)

        # Company name + ticker (large)
        self._text_box(slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.8),
                       f"{name} ({ticker})", size=38, bold=True, color=WHITE)

        # Key metrics row — 4 large cards
        y = Inches(3.0)
        card_w = Inches(2.8)
        gap = Inches(0.25)
        x = Inches(0.8)

        if price is not None:
            change = snap.get("change_1d_pct")
            chg_str = f"{change:+.2f}% today" if change else ""
            self._large_metric_card(slide, x, y, "Current Price",
                                    _fmt_price(price), value_color=WHITE,
                                    description=chg_str)
        x += card_w + gap
        if fair_mid is not None:
            self._large_metric_card(slide, x, y, "Price Target",
                                    _fmt_price(fair_mid), value_color=rec_color,
                                    description=f"Range: {_fmt_price(fair_low)} \u2013 {_fmt_price(fair_high)}" if fair_low and fair_high else "")
        x += card_w + gap
        if upside is not None:
            sign = "+" if upside > 0 else ""
            self._large_metric_card(slide, x, y, "Implied Upside",
                                    f"{sign}{upside:.1f}%",
                                    value_color=ACCENT_GREEN if upside > 0 else ACCENT_RED,
                                    description="Based on blended DCF + comps valuation")
        x += card_w + gap
        # 4th card: Market cap
        mcap = snap.get("market_cap")
        pe = snap.get("pe_ratio")
        if mcap:
            self._large_metric_card(slide, x, y, "Market Cap",
                                    _fmt_large(mcap), value_color=WHITE,
                                    description=f"P/E: {_fmt_ratio(pe)}" if pe else "")

        # Investment thesis paragraph (larger, more detailed)
        thesis_text = thesis.get("investment_thesis", "")
        if thesis_text:
            # Use more of the text (up to ~450 chars)
            display_text = thesis_text if len(thesis_text) <= 450 else thesis_text[:447] + "..."
            self._panel(slide, Inches(0.8), Inches(4.7), Inches(11.6), Inches(1.8),
                        fill=DARK_PANEL, border=BORDER_LIGHT)
            self._text_box(slide, Inches(1.0), Inches(4.8), Inches(1.5), Inches(0.25),
                           "INVESTMENT THESIS", size=8, bold=True, color=ACCENT_BLUE)
            self._text_box(slide, Inches(1.0), Inches(5.1), Inches(11.2), Inches(1.3),
                           display_text, size=13, color=OFF_WHITE, line_spacing=18)

        self._footer(slide, ticker)

    # ------------------------------------------------------------------
    # SLIDE 2: Business Overview
    # ------------------------------------------------------------------

    def _add_business_overview_slide(self, analysis: dict):
        slide = self._blank_slide()
        self._title_bar(slide, "Business Overview")

        snap = analysis.get("snapshot") or {}
        thesis = _sg(analysis, "fundamentals", "investment_thesis") or {}
        fh = _sg(analysis, "fundamentals", "financial_health") or {}
        income = fh.get("income_summary") or {}

        # Left column — company description (expanded)
        overview = thesis.get("company_overview", snap.get("description", ""))
        if overview:
            display_text = overview if len(overview) <= 900 else overview[:897] + "..."
            # Panel background for text
            self._panel(slide, Inches(0.5), Inches(1.15), Inches(7.0), Inches(3.6),
                        fill=DARK_PANEL, border=BORDER_COLOR)
            self._text_box(slide, Inches(0.7), Inches(1.25), Inches(1.5), Inches(0.25),
                           "COMPANY DESCRIPTION", size=8, bold=True, color=ACCENT_BLUE)
            self._text_box(slide, Inches(0.7), Inches(1.55), Inches(6.6), Inches(3.0),
                           display_text, size=12, color=OFF_WHITE, line_spacing=17)

        # Right column — key stats grid (3x3)
        x_base = Inches(7.8)
        y_base = Inches(1.15)
        cw = Inches(1.7)
        ch = Inches(1.05)
        gap_x = Inches(0.1)
        gap_y = Inches(0.1)

        stats = [
            ("Market Cap", _fmt_large(snap.get("market_cap")), WHITE),
            ("Revenue", _fmt_large(snap.get("total_revenue")), WHITE),
            ("Net Income", _fmt_large(income.get("net_income")), WHITE),
            ("Rev Growth", _fmt_pct(snap.get("revenue_growth")), ACCENT_GREEN if _pct_raw(snap.get("revenue_growth")) and _pct_raw(snap.get("revenue_growth")) > 0 else ACCENT_RED),
            ("Op Margin", _fmt_pct(snap.get("operating_margins")), WHITE),
            ("Profit Margin", _fmt_pct(snap.get("profit_margins")), WHITE),
            ("P/E Ratio", _fmt_ratio(snap.get("pe_ratio")), WHITE),
            ("ROE", _fmt_pct(snap.get("return_on_equity")), WHITE),
            ("Employees", f"{snap.get('employees', 0):,}" if snap.get("employees") else "\u2014", WHITE),
        ]

        for i, (label, value, vc) in enumerate(stats):
            row = i // 3
            col = i % 3
            cx = x_base + col * (cw + gap_x)
            cy = y_base + row * (ch + gap_y)
            self._metric_card(slide, cx, cy, label, str(value),
                              width=cw, height=ch, value_color=vc)

        # Bottom section — Key competitive advantages from strengths
        strengths = thesis.get("strengths", [])
        if strengths:
            self._divider(slide, Inches(5.0))
            self._section_header(slide, Inches(0.6), Inches(5.2), "Key Strengths")
            # Show up to 3 strengths with bullet points
            bullet_items = [s for s in strengths[:3]]
            self._bullet_list(slide, Inches(0.8), Inches(5.6), Inches(12), Inches(1.2),
                              bullet_items, size=12, color=OFF_WHITE, spacing=6)

        self._footer(slide, snap.get("ticker", ""))

    # ------------------------------------------------------------------
    # SLIDE 3: Financial Summary
    # ------------------------------------------------------------------

    def _add_financial_summary_slide(self, analysis: dict):
        slide = self._blank_slide()
        self._title_bar(slide, "Financial Summary",
                        subtitle="Key financial metrics from the latest filings")

        fh = _sg(analysis, "fundamentals", "financial_health") or {}
        income = fh.get("income_summary") or {}
        balance = fh.get("balance_summary") or {}
        cashflow = fh.get("cashflow_summary") or {}
        quality = fh.get("quality_scores") or {}
        snap = analysis.get("snapshot") or {}

        col_w = Inches(3.8)
        col_gap = Inches(0.35)
        y_start = Inches(1.3)

        columns = [
            ("Income Statement", [
                ("Revenue", _fmt_large(income.get("revenue"))),
                ("Revenue Growth YoY", _fmt_pct(income.get("revenue_growth_yoy"))),
                ("Gross Margin", _fmt_pct(income.get("gross_margin"))),
                ("Operating Margin", _fmt_pct(income.get("operating_margin"))),
                ("Net Margin", _fmt_pct(income.get("net_margin"))),
                ("EBITDA", _fmt_large(income.get("ebitda"))),
                ("EPS", _fmt_price(income.get("eps"))),
            ]),
            ("Balance Sheet", [
                ("Total Assets", _fmt_large(balance.get("total_assets"))),
                ("Total Equity", _fmt_large(balance.get("total_equity"))),
                ("Total Debt", _fmt_large(balance.get("total_debt"))),
                ("Cash & Equiv", _fmt_large(balance.get("cash"))),
                ("Net Debt", _fmt_large(balance.get("net_debt"))),
                ("Current Ratio", _fmt_ratio(balance.get("current_ratio"))),
                ("Debt / Equity", _fmt_ratio(balance.get("debt_to_equity"))),
            ]),
            ("Cash Flow", [
                ("Operating CF", _fmt_large(cashflow.get("operating_cf"))),
                ("CapEx", _fmt_large(cashflow.get("capex"))),
                ("Free Cash Flow", _fmt_large(cashflow.get("free_cf"))),
                ("FCF Margin", _fmt_pct(cashflow.get("fcf_margin"))),
                ("FCF Yield", _fmt_pct(cashflow.get("fcf_yield"))),
                ("Dividend Yield", _fmt_pct(snap.get("dividend_yield"))),
                ("Book Value/Share", _fmt_price(balance.get("book_value_per_share"))),
            ]),
        ]

        for ci, (col_title, items) in enumerate(columns):
            x = Inches(0.6) + ci * (col_w + col_gap)
            # Column header
            hdr = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, y_start, col_w, Inches(0.4),
            )
            hdr.fill.solid()
            hdr.fill.fore_color.rgb = ACCENT_BLUE
            hdr.line.fill.background()
            self._text_box(slide, x + Inches(0.15), y_start + Inches(0.02),
                           col_w - Inches(0.3), Inches(0.35),
                           col_title, size=11, bold=True, color=WHITE)

            # Rows
            for ri, (label, value) in enumerate(items):
                ry = y_start + Inches(0.5) + ri * Inches(0.65)
                # background stripe
                stripe = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, x, ry, col_w, Inches(0.58),
                )
                stripe.fill.solid()
                stripe.fill.fore_color.rgb = DARK_PANEL if ri % 2 == 0 else NAVY
                stripe.line.fill.background()

                self._text_box(slide, x + Inches(0.15), ry + Inches(0.08),
                               Inches(2.0), Inches(0.42),
                               label, size=10, color=MID_GRAY)
                # Color-code certain values
                v_color = WHITE
                if "growth" in label.lower() or "margin" in label.lower() or "yield" in label.lower():
                    raw = _pct_raw(dict(items).get(label.replace("Revenue Growth YoY", "revenue_growth_yoy")))
                    # Just keep white for simplicity — the value speaks for itself
                self._text_box(slide, x + Inches(2.1), ry + Inches(0.08),
                               Inches(1.5), Inches(0.42),
                               value, size=12, bold=True, color=v_color, align=PP_ALIGN.RIGHT)

        # Quality scores at bottom if available
        accruals = quality.get("accruals_ratio")
        if accruals is not None:
            self._text_box(slide, Inches(0.6), Inches(6.5), Inches(8), Inches(0.3),
                           f"Accruals Ratio: {_fmt_ratio(accruals)}  \u2022  "
                           f"Revenue Growth Std: {_fmt_pct(quality.get('revenue_growth_std'))}",
                           size=9, color=DIM_GRAY)

        self._footer(slide, snap.get("ticker", ""))

    # ------------------------------------------------------------------
    # SLIDE 4: Industry & Competitive Position
    # ------------------------------------------------------------------

    def _add_industry_slide(self, analysis: dict):
        slide = self._blank_slide()
        self._title_bar(slide, "Industry & Competitive Position")

        fund = analysis.get("fundamentals") or {}
        thesis = fund.get("investment_thesis") or {}
        comps = fund.get("comps_valuation") or {}
        snap = analysis.get("snapshot") or {}

        # Left — industry analysis text (expanded with panel)
        ind_text = thesis.get("industry_analysis", "")
        if ind_text:
            display_text = ind_text if len(ind_text) <= 800 else ind_text[:797] + "..."
            self._panel(slide, Inches(0.5), Inches(1.1), Inches(6.5), Inches(3.5),
                        fill=DARK_PANEL, border=BORDER_COLOR)
            self._text_box(slide, Inches(0.7), Inches(1.2), Inches(2), Inches(0.25),
                           "INDUSTRY ANALYSIS", size=8, bold=True, color=ACCENT_BLUE)
            self._text_box(slide, Inches(0.7), Inches(1.5), Inches(6.1), Inches(2.9),
                           display_text, size=12, color=OFF_WHITE, line_spacing=17)
        else:
            # Fallback: show sector/industry info
            sector = snap.get("sector", "")
            industry = snap.get("industry", "")
            if sector or industry:
                self._panel(slide, Inches(0.5), Inches(1.1), Inches(6.5), Inches(2.0),
                            fill=DARK_PANEL, border=BORDER_COLOR)
                self._text_box(slide, Inches(0.7), Inches(1.3), Inches(6.1), Inches(0.5),
                               f"Sector: {sector}", size=16, bold=True, color=WHITE)
                self._text_box(slide, Inches(0.7), Inches(1.9), Inches(6.1), Inches(0.5),
                               f"Industry: {industry}", size=14, color=LIGHT_GRAY)

        # Right — peer comparison chart
        chart_path = self._create_peer_chart(comps, snap)
        if chart_path:
            self._add_picture(slide, chart_path, Inches(7.2), Inches(1.0),
                              width=Inches(5.8), height=Inches(3.5))

        # Bottom — percentile rank badges (wider, with labels)
        pct_rank = comps.get("percentile_rank") or {}
        if pct_rank:
            self._divider(slide, Inches(4.8))
            self._section_header(slide, Inches(0.6), Inches(4.95), "Percentile Rank vs Peers")
            bx = Inches(0.6)
            by = Inches(5.4)
            badge_w = Inches(1.5)
            badge_h = Inches(0.85)
            for metric, pctile in list(pct_rank.items())[:8]:
                try:
                    pval = float(pctile)
                except (TypeError, ValueError):
                    continue
                color = ACCENT_GREEN if pval >= 60 else (ACCENT_RED if pval <= 30 else LIGHT_GRAY)
                label_clean = metric.replace("_", " ").title()
                # Contextual subtitle
                ctx = "Top quartile" if pval >= 75 else ("Above median" if pval >= 50 else ("Below median" if pval >= 25 else "Bottom quartile"))
                self._metric_card(slide, bx, by, label_clean, f"{pval:.0f}th",
                                  width=badge_w, height=badge_h, value_color=color,
                                  subtitle=ctx, subtitle_color=DIM_GRAY)
                bx += badge_w + Inches(0.1)

        self._footer(slide, snap.get("ticker", ""))

    # ------------------------------------------------------------------
    # SLIDE 5: Why Mispriced?
    # ------------------------------------------------------------------

    def _add_why_mispriced_slide(self, analysis: dict):
        slide = self._blank_slide()
        self._title_bar(slide, "Why Is This Mispriced?",
                        subtitle="Identifying the gap between market price and intrinsic value")

        fund = analysis.get("fundamentals") or {}
        verdict = fund.get("verdict") or {}
        news = _sg(analysis, "news", "synthesis") or {}
        thesis = fund.get("investment_thesis") or {}
        combined = fund.get("combined_range") or {}
        snap = analysis.get("snapshot") or {}

        price = snap.get("price")
        low = combined.get("low")
        mid = combined.get("mid")
        high = combined.get("high")

        # Fair value range chart (full width)
        chart_path = self._create_fair_value_chart(price, low, mid, high)
        if chart_path:
            self._add_picture(slide, chart_path, Inches(0.4), Inches(1.1),
                              width=Inches(12.5), height=Inches(2.2))

        # Two-column layout below chart
        # Left column — Market Sentiment (enriched)
        y = Inches(3.5)
        self._panel(slide, Inches(0.5), y, Inches(6.0), Inches(3.1),
                    fill=DARK_PANEL, border=BORDER_COLOR)

        self._text_box(slide, Inches(0.7), y + Inches(0.1), Inches(3), Inches(0.25),
                       "MARKET SENTIMENT", size=8, bold=True, color=ACCENT_BLUE)

        sentiment = news.get("overall_sentiment", "")
        score = news.get("sentiment_score")
        momentum = news.get("news_momentum", "")
        narrative = news.get("narrative", "")
        agg = news.get("aggregate_sentiment") or {}
        pos_pct = agg.get("positive_pct")
        neg_pct = agg.get("negative_pct")
        article_count = news.get("article_count")

        if sentiment:
            s_color = ACCENT_GREEN if sentiment == "bullish" else (
                ACCENT_RED if sentiment == "bearish" else LIGHT_GRAY)
            # Sentiment + score in one line
            score_str = f"  (score: {score:+.2f})" if score is not None else ""
            self._text_box(slide, Inches(0.7), y + Inches(0.4), Inches(5.5), Inches(0.45),
                           f"{sentiment.upper()}{score_str}", size=20, bold=True, color=s_color)

        # Momentum + article count
        meta_parts = []
        if momentum:
            meta_parts.append(f"Momentum: {momentum.capitalize()}")
        if article_count:
            meta_parts.append(f"{article_count} articles analyzed")
        if pos_pct is not None and neg_pct is not None:
            meta_parts.append(f"{_fmt_pct(pos_pct)} positive / {_fmt_pct(neg_pct)} negative")
        if meta_parts:
            self._text_box(slide, Inches(0.7), y + Inches(0.9), Inches(5.5), Inches(0.35),
                           "  \u2022  ".join(meta_parts), size=10, color=MID_GRAY)

        if narrative:
            trunc = narrative if len(narrative) <= 350 else narrative[:347] + "..."
            self._text_box(slide, Inches(0.7), y + Inches(1.35), Inches(5.5), Inches(1.6),
                           trunc, size=11, color=OFF_WHITE, line_spacing=16)

        # Right column — Catalysts (enriched)
        catalysts = thesis.get("catalysts", [])
        news_catalysts = news.get("catalysts") or []

        self._panel(slide, Inches(6.7), y, Inches(6.0), Inches(3.1),
                    fill=DARK_PANEL, border=BORDER_COLOR)
        self._text_box(slide, Inches(6.9), y + Inches(0.1), Inches(3), Inches(0.25),
                       "KEY CATALYSTS", size=8, bold=True, color=ACCENT_GOLD)

        all_catalysts = catalysts[:4]
        # Add news-derived catalysts if thesis ones are short
        if len(all_catalysts) < 4 and news_catalysts:
            for nc in news_catalysts:
                if len(all_catalysts) >= 5:
                    break
                if isinstance(nc, dict):
                    event = nc.get("event", nc.get("headline", ""))
                    if event and event not in all_catalysts:
                        all_catalysts.append(event)
                elif isinstance(nc, str) and nc not in all_catalysts:
                    all_catalysts.append(nc)

        if all_catalysts:
            cat_paras = []
            for ci, cat in enumerate(all_catalysts[:5]):
                cat_paras.append({
                    "text": f"{ci + 1}.  {cat}",
                    "size": 12,
                    "color": OFF_WHITE,
                    "spacing_after": 10,
                })
            self._rich_text_box(slide, Inches(6.9), y + Inches(0.45),
                                Inches(5.6), Inches(2.5), cat_paras)

        # Verdict reasoning at very bottom
        reasoning = verdict.get("reasoning", "")
        if reasoning:
            self._text_box(slide, Inches(0.6), Inches(6.55), Inches(12), Inches(0.35),
                           reasoning[:300], size=9, color=DIM_GRAY)

        self._footer(slide, snap.get("ticker", ""))

    # ------------------------------------------------------------------
    # SLIDES 6-8: Core Theses
    # ------------------------------------------------------------------

    def _add_thesis_slides(self, analysis: dict):
        fund = analysis.get("fundamentals") or {}
        thesis = fund.get("investment_thesis") or {}
        snap = analysis.get("snapshot") or {}
        forecast_data = analysis.get("forecast") or {}
        verdict = fund.get("verdict") or {}

        strengths = thesis.get("strengths", [])
        catalysts = list(thesis.get("catalysts", []))  # copy
        risks = thesis.get("risks", [])
        key_metrics = thesis.get("key_metrics", [])

        # Build thesis items: combine strengths + catalysts
        items: list[tuple[str, str | None, str | None]] = []
        for s in strengths[:3]:
            related_cat = catalysts.pop(0) if catalysts else None
            related_risk = risks[len(items)] if len(items) < len(risks) else None
            items.append((s, related_cat, related_risk))

        # Ensure at least 1 slide
        if not items:
            items = [("Investment opportunity based on fundamental analysis.", None, None)]

        for i, (strength, catalyst, risk) in enumerate(items[:3]):
            slide = self._blank_slide()
            self._title_bar(slide, f"Core Thesis {i + 1}")

            # Large thesis number accent
            self._text_box(slide, Inches(0.6), Inches(1.1), Inches(1), Inches(1),
                           str(i + 1), size=60, bold=True, color=ACCENT_BLUE_DIM)

            # Main thesis text (larger, bolder)
            self._text_box(slide, Inches(1.5), Inches(1.2), Inches(11), Inches(1.4),
                           strength, size=20, bold=True, color=WHITE, line_spacing=28)

            # Catalyst section with panel
            cy = Inches(2.9)
            if catalyst:
                self._panel(slide, Inches(0.5), cy, Inches(12.3), Inches(1.5),
                            fill=DARK_PANEL, border=BORDER_COLOR)
                self._side_accent(slide, cy + Inches(0.1), Inches(1.3), color=ACCENT_GOLD)
                self._text_box(slide, Inches(0.85), cy + Inches(0.1), Inches(3), Inches(0.25),
                               "CATALYST", size=9, bold=True, color=ACCENT_GOLD)
                self._text_box(slide, Inches(0.85), cy + Inches(0.4), Inches(11.7), Inches(1.0),
                               catalyst, size=14, color=OFF_WHITE, line_spacing=20)
                cy += Inches(1.7)

            # Counter-argument / risk awareness
            if risk and i < 2:  # Show risk on first 2 thesis slides
                self._panel(slide, Inches(0.5), cy, Inches(12.3), Inches(1.0),
                            fill=DARK_PANEL, border=BORDER_COLOR)
                self._side_accent(slide, cy + Inches(0.1), Inches(0.8), color=ACCENT_RED)
                self._text_box(slide, Inches(0.85), cy + Inches(0.1), Inches(3), Inches(0.25),
                               "KEY RISK TO THESIS", size=9, bold=True, color=ACCENT_RED)
                self._text_box(slide, Inches(0.85), cy + Inches(0.4), Inches(11.7), Inches(0.5),
                               risk[:200], size=12, color=LIGHT_GRAY)
                cy += Inches(1.2)

            # Supporting metrics (4 cards)
            met_y = max(cy, Inches(4.8))
            # Get metrics relevant to this thesis (cycle through)
            start_idx = i * 3
            slide_metrics = key_metrics[start_idx:start_idx + 4]
            if not slide_metrics:
                slide_metrics = key_metrics[:4]

            if slide_metrics:
                self._section_header(slide, Inches(0.6), met_y - Inches(0.35), "Supporting Data")
                for mi, km in enumerate(slide_metrics[:4]):
                    label = km.get("label", "")
                    value = km.get("value", "")
                    context = km.get("context", "")
                    if label and value:
                        self._metric_card(
                            slide, Inches(0.6) + mi * Inches(3.1), met_y,
                            label, str(value), width=Inches(2.9), height=Inches(1.0),
                            subtitle=context[:50] if context else "",
                            subtitle_color=DIM_GRAY,
                        )

            # Forecast on first thesis slide
            if i == 0:
                horizons = forecast_data.get("forecasts") or {}
                horizon_12m = horizons.get("12M") or {}
                ep = horizon_12m.get("ensemble_prediction") or {}
                pred_ret = ep.get("predicted_return")
                confidence = ep.get("confidence")
                direction = ep.get("direction", "")
                if pred_ret is not None:
                    parts = [f"12-Month ML Forecast: {pred_ret:+.1f}% return"]
                    if confidence:
                        parts.append(f"{confidence * 100:.0f}% confidence")
                    if direction:
                        parts.append(f"direction: {direction}")
                    self._text_box(slide, Inches(0.6), Inches(6.3), Inches(10), Inches(0.35),
                                   "  \u2022  ".join(parts),
                                   size=11, bold=True, color=ACCENT_BLUE)

            self._footer(slide, snap.get("ticker", ""))

    # ------------------------------------------------------------------
    # SLIDE 9: Downside Protection
    # ------------------------------------------------------------------

    def _add_downside_slide(self, analysis: dict):
        slide = self._blank_slide()
        self._title_bar(slide, "Downside Protection",
                        subtitle="Risk management and portfolio resilience analysis")

        risk = analysis.get("risk") or {}
        var_data = risk.get("var") or {}
        dd_data = risk.get("drawdown") or {}
        scenarios = risk.get("scenarios") or []
        fund = analysis.get("fundamentals") or {}
        fh = fund.get("financial_health") or {}
        cashflow = fh.get("cashflow_summary") or {}
        balance = fh.get("balance_summary") or {}
        quant = analysis.get("quant") or {}
        perf = quant.get("performance") or {}
        snap = analysis.get("snapshot") or {}

        # Top row: Three panels
        pw = Inches(3.9)
        ph = Inches(2.8)
        pg = Inches(0.2)
        py = Inches(1.2)

        # Panel 1: Value at Risk
        self._panel(slide, Inches(0.5), py, pw, ph, fill=DARK_PANEL, border=BORDER_COLOR)
        self._text_box(slide, Inches(0.7), py + Inches(0.1), Inches(3), Inches(0.25),
                       "VALUE AT RISK (95%, 1-DAY)", size=8, bold=True, color=ACCENT_BLUE)

        var_items = [
            ("Parametric VaR", var_data.get("parametric_var_pct"), "Normal distribution model"),
            ("Historical VaR", var_data.get("historical_var_pct"), "Based on actual returns"),
            ("Monte Carlo VaR", var_data.get("monte_carlo_var_pct"), "10,000 simulated paths"),
            ("CVaR (Exp. Shortfall)", var_data.get("expected_shortfall_pct"), "Expected loss beyond VaR"),
        ]
        for vi, (label, value, desc) in enumerate(var_items):
            vy = py + Inches(0.45) + vi * Inches(0.55)
            self._text_box(slide, Inches(0.7), vy, Inches(2.0), Inches(0.25),
                           label, size=10, color=LIGHT_GRAY)
            self._text_box(slide, Inches(2.8), vy, Inches(1.2), Inches(0.25),
                           _fmt_pct(value), size=11, bold=True, color=ACCENT_RED,
                           align=PP_ALIGN.RIGHT)
            self._text_box(slide, Inches(0.7), vy + Inches(0.22), Inches(3), Inches(0.18),
                           desc, size=7, color=DIM_GRAY)

        # Daily & Annualized Vol
        daily_vol = var_data.get("daily_vol")
        ann_vol = var_data.get("annualized_vol")
        if daily_vol or ann_vol:
            self._text_box(slide, Inches(0.7), py + Inches(2.45), Inches(3.2), Inches(0.25),
                           f"Daily Vol: {_fmt_pct(daily_vol)}  \u2022  Ann. Vol: {_fmt_pct(ann_vol)}",
                           size=9, color=MID_GRAY)

        # Panel 2: Cash Flow Resilience + Drawdown
        p2x = Inches(0.5) + pw + pg
        self._panel(slide, p2x, py, pw, ph, fill=DARK_PANEL, border=BORDER_COLOR)
        self._text_box(slide, p2x + Inches(0.2), py + Inches(0.1), Inches(3), Inches(0.25),
                       "BALANCE SHEET & CASH FLOW", size=8, bold=True, color=ACCENT_BLUE)

        cash_items = [
            ("Free Cash Flow", _fmt_large(cashflow.get("free_cf"))),
            ("FCF Margin", _fmt_pct(cashflow.get("fcf_margin"))),
            ("Cash Position", _fmt_large(balance.get("cash"))),
            ("Current Ratio", _fmt_ratio(balance.get("current_ratio"))),
            ("Debt / Equity", _fmt_ratio(balance.get("debt_to_equity"))),
        ]
        for vi, (label, value) in enumerate(cash_items):
            vy = py + Inches(0.45) + vi * Inches(0.45)
            self._text_box(slide, p2x + Inches(0.2), vy, Inches(2.2), Inches(0.35),
                           label, size=10, color=LIGHT_GRAY)
            self._text_box(slide, p2x + Inches(2.5), vy, Inches(1.2), Inches(0.35),
                           value, size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

        # Panel 3: Risk-Adjusted Returns + Drawdown
        p3x = Inches(0.5) + (pw + pg) * 2
        self._panel(slide, p3x, py, pw, ph, fill=DARK_PANEL, border=BORDER_COLOR)
        self._text_box(slide, p3x + Inches(0.2), py + Inches(0.1), Inches(3), Inches(0.25),
                       "RISK-ADJUSTED RETURNS", size=8, bold=True, color=ACCENT_BLUE)

        risk_items = [
            ("Sharpe Ratio", _fmt_ratio(perf.get("sharpe_ratio")),
             "Risk-adjusted excess return"),
            ("Sortino Ratio", _fmt_ratio(perf.get("sortino_ratio")),
             "Downside deviation adjusted"),
            ("Max Drawdown", _fmt_pct(perf.get("max_drawdown")),
             "Largest peak-to-trough decline"),
            ("Ann. Volatility", _fmt_pct(perf.get("annualized_vol")),
             "Annualized std deviation"),
        ]
        for vi, (label, value, desc) in enumerate(risk_items):
            vy = py + Inches(0.45) + vi * Inches(0.55)
            self._text_box(slide, p3x + Inches(0.2), vy, Inches(2.0), Inches(0.25),
                           label, size=10, color=LIGHT_GRAY)
            self._text_box(slide, p3x + Inches(2.3), vy, Inches(1.3), Inches(0.25),
                           value, size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
            self._text_box(slide, p3x + Inches(0.2), vy + Inches(0.22), Inches(3), Inches(0.18),
                           desc, size=7, color=DIM_GRAY)

        # Bottom: Historical stress tests
        self._divider(slide, Inches(4.2))
        self._section_header(slide, Inches(0.6), Inches(4.35), "Historical Stress Tests")

        if isinstance(scenarios, list) and scenarios:
            # Header row
            headers = ["Scenario", "Stock Return", "Benchmark", "Max Drawdown", "Period"]
            hx_positions = [Inches(0.6), Inches(5), Inches(7.2), Inches(9.4), Inches(11.2)]
            hx_widths = [Inches(4.2), Inches(2), Inches(2), Inches(1.8), Inches(2)]
            for hi, header in enumerate(headers):
                self._text_box(slide, hx_positions[hi], Inches(4.75),
                               hx_widths[hi], Inches(0.25),
                               header, size=8, bold=True, color=MID_GRAY)

            for si, sc in enumerate(scenarios[:4]):
                if not isinstance(sc, dict):
                    continue
                sy = Inches(5.05) + si * Inches(0.42)
                # Alternating stripe
                stripe = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0.5), sy, Inches(12.33), Inches(0.38),
                )
                stripe.fill.solid()
                stripe.fill.fore_color.rgb = DARK_PANEL if si % 2 == 0 else NAVY
                stripe.line.fill.background()

                name = sc.get("description", sc.get("name", sc.get("scenario", f"Scenario {si + 1}")))
                stock_ret = sc.get("stock_return") or sc.get("return") or sc.get("beta_adjusted_return")
                bench_ret = sc.get("benchmark_return")
                mdd = sc.get("max_drawdown")
                period = sc.get("period", "")

                self._text_box(slide, hx_positions[0], sy + Inches(0.04), hx_widths[0], Inches(0.3),
                               str(name)[:45], size=10, color=OFF_WHITE)

                ret_color = ACCENT_RED
                if stock_ret is not None:
                    try:
                        if float(stock_ret) > 0:
                            ret_color = ACCENT_GREEN
                    except (TypeError, ValueError):
                        pass
                self._text_box(slide, hx_positions[1], sy + Inches(0.04), hx_widths[1], Inches(0.3),
                               _fmt_pct(stock_ret), size=10, bold=True, color=ret_color)
                self._text_box(slide, hx_positions[2], sy + Inches(0.04), hx_widths[2], Inches(0.3),
                               _fmt_pct(bench_ret), size=10, color=LIGHT_GRAY)
                self._text_box(slide, hx_positions[3], sy + Inches(0.04), hx_widths[3], Inches(0.3),
                               _fmt_pct(mdd), size=10, color=LIGHT_GRAY)
                self._text_box(slide, hx_positions[4], sy + Inches(0.04), hx_widths[4], Inches(0.3),
                               str(period)[:20], size=9, color=DIM_GRAY)

        self._footer(slide, snap.get("ticker", ""))

    # ------------------------------------------------------------------
    # SLIDE 10: Valuation
    # ------------------------------------------------------------------

    def _add_valuation_slide(self, analysis: dict):
        slide = self._blank_slide()
        self._title_bar(slide, "Valuation",
                        subtitle="DCF and comparable company analysis")

        fund = analysis.get("fundamentals") or {}
        dcf = fund.get("dcf_valuation") or {}
        comps = fund.get("comps_valuation") or {}
        combined = fund.get("combined_range") or {}
        snap = analysis.get("snapshot") or {}
        price = snap.get("price")

        # Top-left: DCF section with panel
        dcf_w = Inches(6.2)
        self._panel(slide, Inches(0.5), Inches(1.15), dcf_w, Inches(2.2),
                    fill=DARK_PANEL, border=BORDER_COLOR)
        self._text_box(slide, Inches(0.7), Inches(1.25), Inches(3), Inches(0.25),
                       "DCF VALUATION", size=8, bold=True, color=ACCENT_BLUE)

        # DCF hero metrics
        dcf_cards = [
            ("Implied Price", _fmt_price(dcf.get("implied_price")), ACCENT_GREEN),
            ("Enterprise Value", _fmt_large(dcf.get("enterprise_value")), WHITE),
            ("Equity Value", _fmt_large(dcf.get("equity_value")), WHITE),
        ]
        for di, (label, value, vc) in enumerate(dcf_cards):
            self._metric_card(slide, Inches(0.7) + di * Inches(1.85), Inches(1.6),
                              label, value, width=Inches(1.75), height=Inches(0.8),
                              value_color=vc)

        # DCF assumptions (more detailed)
        assumptions = dcf.get("assumptions_used") or {}
        if assumptions:
            self._text_box(slide, Inches(0.7), Inches(2.55), Inches(5.8), Inches(0.25),
                           "ASSUMPTIONS", size=7, bold=True, color=MID_GRAY)
            a_parts = []
            wacc = assumptions.get("wacc")
            tg = assumptions.get("terminal_growth_rate")
            om = assumptions.get("operating_margin")
            tr = assumptions.get("tax_rate")
            if wacc:
                a_parts.append(f"WACC: {_fmt_pct(wacc)}")
            if tg:
                a_parts.append(f"Terminal Growth: {_fmt_pct(tg)}")
            if om:
                a_parts.append(f"Op Margin: {_fmt_pct(om)}")
            if tr:
                a_parts.append(f"Tax Rate: {_fmt_pct(tr)}")
            proj_years = assumptions.get("projection_years")
            if proj_years:
                a_parts.append(f"{proj_years}yr projection")
            self._text_box(slide, Inches(0.7), Inches(2.78), Inches(5.8), Inches(0.4),
                           "   \u2022   ".join(a_parts), size=9, color=LIGHT_GRAY)

        # Upside/downside callout
        dcf_upside = dcf.get("upside_downside_pct")
        if dcf_upside is not None:
            try:
                uv = float(dcf_upside)
                uc = ACCENT_GREEN if uv > 0 else ACCENT_RED
                sign = "+" if uv > 0 else ""
                self._text_box(slide, Inches(0.7), Inches(3.1), Inches(5.8), Inches(0.25),
                               f"DCF implies {sign}{uv:.1f}% {'upside' if uv > 0 else 'downside'} from current price",
                               size=10, bold=True, color=uc)
            except (TypeError, ValueError):
                pass

        # Top-right: Comps section with panel
        comp_x = Inches(6.9)
        comp_w = Inches(5.9)
        self._panel(slide, comp_x, Inches(1.15), comp_w, Inches(2.2),
                    fill=DARK_PANEL, border=BORDER_COLOR)
        self._text_box(slide, comp_x + Inches(0.2), Inches(1.25), Inches(4), Inches(0.25),
                       "COMPARABLE COMPANY ANALYSIS", size=8, bold=True, color=ACCENT_BLUE)

        val_range = comps.get("valuation_range") or {}
        approaches = [
            ("P/E Approach", val_range.get("pe_approach") or {}),
            ("EV/EBITDA Approach", val_range.get("ev_ebitda_approach") or {}),
            ("EV/Revenue Approach", val_range.get("ev_revenue_approach") or {}),
        ]

        for ai, (approach_name, approach_data) in enumerate(approaches):
            ay = Inches(1.65) + ai * Inches(0.55)
            implied = approach_data.get("implied_value") or approach_data.get("implied_market_cap")
            upside = approach_data.get("implied_upside")
            peer_med = approach_data.get("peer_median_pe") or approach_data.get("peer_median_ev_ebitda") or approach_data.get("peer_median_ev_revenue")

            self._text_box(slide, comp_x + Inches(0.2), ay, Inches(2.2), Inches(0.25),
                           approach_name, size=10, color=LIGHT_GRAY)
            self._text_box(slide, comp_x + Inches(2.5), ay, Inches(1.5), Inches(0.25),
                           _fmt_large(implied), size=11, bold=True, color=WHITE,
                           align=PP_ALIGN.RIGHT)
            if upside is not None:
                try:
                    uv = float(upside)
                    uc = ACCENT_GREEN if uv > 0 else ACCENT_RED
                    self._text_box(slide, comp_x + Inches(4.2), ay, Inches(1.2), Inches(0.25),
                                   f"{uv:+.1f}%", size=10, bold=True, color=uc,
                                   align=PP_ALIGN.RIGHT)
                except (TypeError, ValueError):
                    pass

        # Peer count
        peer_count = comps.get("peer_count")
        if peer_count:
            self._text_box(slide, comp_x + Inches(0.2), Inches(3.1), Inches(5), Inches(0.2),
                           f"Based on {peer_count} comparable companies",
                           size=9, color=DIM_GRAY)

        # Bottom: Sensitivity heatmap + Bull/Base/Bear
        self._divider(slide, Inches(3.5))

        # Sensitivity heatmap
        heatmap_path = self._create_sensitivity_heatmap(dcf, price)
        if heatmap_path:
            self._add_picture(slide, heatmap_path, Inches(0.3), Inches(3.7),
                              width=Inches(6.8), height=Inches(3.3))

        # Bull / Base / Bear range chart
        range_path = self._create_range_chart(combined, price)
        if range_path:
            self._add_picture(slide, range_path, Inches(7.2), Inches(3.7),
                              width=Inches(5.5), height=Inches(1.6))

        # Combined range metrics
        if combined:
            y_range = Inches(5.6)
            mw = combined.get("methodology_weights") or {}
            dcf_weight = mw.get("dcf")
            comps_weight = mw.get("comps")

            self._large_metric_card(slide, Inches(7.2), y_range, "Bear Case",
                                    _fmt_price(combined.get("low")),
                                    width=Inches(1.7), height=Inches(1.2), value_color=ACCENT_RED,
                                    description="Downside scenario")
            self._large_metric_card(slide, Inches(9.1), y_range, "Base Case",
                                    _fmt_price(combined.get("mid")),
                                    width=Inches(1.7), height=Inches(1.2), value_color=ACCENT_BLUE,
                                    description="Most likely")
            self._large_metric_card(slide, Inches(11.0), y_range, "Bull Case",
                                    _fmt_price(combined.get("high")),
                                    width=Inches(1.7), height=Inches(1.2), value_color=ACCENT_GREEN,
                                    description="Upside scenario")

            if dcf_weight and comps_weight:
                self._text_box(slide, Inches(7.2), Inches(6.5), Inches(5.5), Inches(0.2),
                               f"Blended: {dcf_weight*100:.0f}% DCF + {comps_weight*100:.0f}% Comps",
                               size=8, color=DIM_GRAY, align=PP_ALIGN.CENTER)

        self._footer(slide, snap.get("ticker", ""))

    # ------------------------------------------------------------------
    # SLIDE 11: Risks & Mitigants
    # ------------------------------------------------------------------

    def _add_risks_slide(self, analysis: dict):
        slide = self._blank_slide()
        self._title_bar(slide, "Risks & Mitigants",
                        subtitle="Key risks to the investment thesis and how they are addressed")

        fund = analysis.get("fundamentals") or {}
        thesis = fund.get("investment_thesis") or {}
        news = _sg(analysis, "news", "synthesis") or {}
        snap = analysis.get("snapshot") or {}

        risks = thesis.get("risks", [])
        strengths = thesis.get("strengths", [])
        news_risks = news.get("risks") or []

        # Supplement risks from news if needed
        all_risks = list(risks[:5])
        if len(all_risks) < 3 and news_risks:
            for nr in news_risks:
                if len(all_risks) >= 5:
                    break
                if isinstance(nr, dict):
                    event = nr.get("event", nr.get("headline", ""))
                    if event and event not in all_risks:
                        all_risks.append(event)

        # Build risk/mitigant pairs
        pairs: list[tuple[str, str]] = []
        for i, risk_text in enumerate(all_risks[:5]):
            mitigant = strengths[i] if i < len(strengths) else "Fundamental value and diversified revenue streams provide margin of safety."
            pairs.append((risk_text, mitigant))

        if not pairs:
            self._text_box(slide, Inches(0.6), Inches(2), Inches(10), Inches(1),
                           "No significant risks identified in the current analysis. "
                           "This may indicate limited data availability rather than an absence of risk.",
                           size=14, color=LIGHT_GRAY)
            self._footer(slide, snap.get("ticker", ""))
            return

        # Headers
        self._panel(slide, Inches(0.5), Inches(1.25), Inches(5.9), Inches(0.4),
                    fill=ACCENT_RED, border=ACCENT_RED, radius=False)
        self._text_box(slide, Inches(0.7), Inches(1.28), Inches(5.5), Inches(0.35),
                       "RISK", size=11, bold=True, color=WHITE)

        self._panel(slide, Inches(6.9), Inches(1.25), Inches(5.9), Inches(0.4),
                    fill=ACCENT_GREEN_DIM, border=ACCENT_GREEN_DIM, radius=False)
        self._text_box(slide, Inches(7.1), Inches(1.28), Inches(5.5), Inches(0.35),
                       "MITIGANT", size=11, bold=True, color=WHITE)

        row_h = Inches(1.0) if len(pairs) <= 4 else Inches(0.85)

        for pi, (risk_text, mitigant) in enumerate(pairs):
            y = Inches(1.8) + pi * row_h

            # Row backgrounds
            bg_color = DARK_PANEL if pi % 2 == 0 else NAVY
            self._panel(slide, Inches(0.5), y, Inches(5.9), row_h - Inches(0.08),
                        fill=bg_color, border=BORDER_COLOR, radius=False)
            self._panel(slide, Inches(6.9), y, Inches(5.9), row_h - Inches(0.08),
                        fill=bg_color, border=BORDER_COLOR, radius=False)

            # Risk number badge
            num_badge = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(0.65), y + Inches(0.12), Inches(0.28), Inches(0.28),
            )
            num_badge.fill.solid()
            num_badge.fill.fore_color.rgb = ACCENT_RED
            num_badge.line.fill.background()
            self._text_box(slide, Inches(0.65), y + Inches(0.12), Inches(0.28), Inches(0.28),
                           str(pi + 1), size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

            # Risk text (larger, more readable)
            r_text = risk_text if len(risk_text) <= 200 else risk_text[:197] + "..."
            self._text_box(slide, Inches(1.05), y + Inches(0.1), Inches(5.15), row_h - Inches(0.2),
                           r_text, size=11, color=OFF_WHITE, line_spacing=15)

            # Arrow connector
            self._text_box(slide, Inches(6.45), y + Inches(0.15), Inches(0.4), Inches(0.5),
                           "\u2192", size=18, color=MID_GRAY, align=PP_ALIGN.CENTER)

            # Mitigant text
            m_text = mitigant if len(mitigant) <= 200 else mitigant[:197] + "..."
            self._text_box(slide, Inches(7.1), y + Inches(0.1), Inches(5.5), row_h - Inches(0.2),
                           m_text, size=11, color=OFF_WHITE, line_spacing=15)

        self._footer(slide, snap.get("ticker", ""))

    # ------------------------------------------------------------------
    # SLIDE 12: Conclusion
    # ------------------------------------------------------------------

    def _add_conclusion_slide(self, analysis: dict):
        slide = self._blank_slide()

        snap = analysis.get("snapshot") or {}
        fund = analysis.get("fundamentals") or {}
        verdict = fund.get("verdict") or {}
        thesis = fund.get("investment_thesis") or {}
        forecast = analysis.get("forecast") or {}
        combined = fund.get("combined_range") or {}
        risk_data = analysis.get("risk") or {}
        quant = analysis.get("quant") or {}
        perf = quant.get("performance") or {}

        ticker = snap.get("ticker", analysis.get("ticker", ""))
        name = snap.get("name", ticker)
        label = verdict.get("label", "")
        upside = verdict.get("upside_pct")
        fair_mid = verdict.get("fair_value_mid")
        fair_low = combined.get("low")
        fair_high = combined.get("high")
        price = snap.get("price")

        is_long = label in ("undervalued", "fairly_valued")
        rec_text = "LONG" if is_long else "SHORT"
        rec_color = ACCENT_GREEN if is_long else ACCENT_RED

        # Title
        self._text_box(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.7),
                       "Conclusion", size=34, bold=True, color=WHITE)

        # Recommendation badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(1.5), Inches(2.0), Inches(0.6),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = rec_color
        badge.line.fill.background()
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = rec_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

        # Company + target line
        target_parts = [f"{name} ({ticker})"]
        if fair_mid:
            target_parts.append(f"Target: {_fmt_price(fair_mid)}")
        if upside:
            target_parts.append(f"Upside: {upside:+.1f}%")
        self._text_box(slide, Inches(3.1), Inches(1.5), Inches(9.5), Inches(0.6),
                       "  \u2022  ".join(target_parts),
                       size=18, bold=True, color=WHITE)

        # Key metrics summary row
        my = Inches(2.4)
        summary_metrics = []
        if price:
            summary_metrics.append(("Price", _fmt_price(price), WHITE))
        if fair_low and fair_high:
            summary_metrics.append(("Fair Value Range", f"{_fmt_price(fair_low)} \u2013 {_fmt_price(fair_high)}", ACCENT_BLUE))
        sharpe = perf.get("sharpe_ratio")
        if sharpe:
            summary_metrics.append(("Sharpe", _fmt_ratio(sharpe), WHITE))
        mcap = snap.get("market_cap")
        if mcap:
            summary_metrics.append(("Mkt Cap", _fmt_large(mcap), WHITE))

        for mi, (ml, mv, mc) in enumerate(summary_metrics[:4]):
            self._metric_card(slide, Inches(0.8) + mi * Inches(3.1), my,
                              ml, mv, width=Inches(2.9), height=Inches(0.85),
                              value_color=mc)

        # Key investment points (expanded, numbered)
        points = []
        strengths = thesis.get("strengths", [])
        catalysts = thesis.get("catalysts", [])
        reasoning = verdict.get("reasoning", "")

        if strengths:
            points.append(strengths[0])
        if catalysts:
            points.append(f"Catalyst: {catalysts[0]}")
        if reasoning:
            points.append(reasoning)

        overall_dir = forecast.get("overall_direction", "")
        lt_outlook = forecast.get("long_term_outlook", "")
        if overall_dir:
            parts = [f"ML forecast direction: {overall_dir}"]
            if lt_outlook:
                parts.append(f"long-term outlook: {lt_outlook}")
            points.append(". ".join(parts) + ".")

        # Add a risk awareness point
        risks = thesis.get("risks", [])
        if risks:
            points.append(f"Key risk: {risks[0]}")

        if points:
            self._panel(slide, Inches(0.7), Inches(3.5), Inches(11.9), Inches(2.8),
                        fill=DARK_PANEL, border=BORDER_LIGHT)
            self._text_box(slide, Inches(0.9), Inches(3.6), Inches(3), Inches(0.25),
                           "KEY TAKEAWAYS", size=8, bold=True, color=ACCENT_BLUE)

            point_paras = []
            for idx, pt in enumerate(points[:5]):
                display = pt if len(pt) <= 200 else pt[:197] + "..."
                point_paras.append({
                    "text": f"{idx + 1}.  {display}",
                    "size": 13,
                    "color": OFF_WHITE,
                    "spacing_after": 10,
                })
            self._rich_text_box(slide, Inches(1.0), Inches(3.9),
                                Inches(11.4), Inches(2.3), point_paras)

        # Bottom branding
        self._text_box(slide, Inches(0.8), Inches(6.3), Inches(12), Inches(0.45),
                       "Generated by AlphaEdge Analysis Platform",
                       size=13, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

        self._footer(slide, ticker)

    # ------------------------------------------------------------------
    # Chart generators (matplotlib -> PNG)
    # ------------------------------------------------------------------

    def _create_peer_chart(self, comps: dict, snap: dict) -> str | None:
        """Horizontal bar chart: target vs peer median for key metrics."""
        target = comps.get("target") or {}
        peers = comps.get("peers") or []
        if not target or not peers:
            return None

        metrics = ["pe", "ev_ebitda", "operating_margin", "roe", "revenue_growth", "gross_margin"]
        labels = ["P/E", "EV/EBITDA", "Op Margin", "ROE", "Rev Growth", "Gross Margin"]

        target_vals = []
        peer_medians = []
        valid_labels = []

        for m, l in zip(metrics, labels):
            tv = target.get(m)
            peer_vs = [p.get(m) for p in peers if p.get(m) is not None]
            if tv is not None and peer_vs:
                try:
                    target_vals.append(float(tv))
                    peer_medians.append(float(np.median(peer_vs)))
                    valid_labels.append(l)
                except (TypeError, ValueError):
                    pass

        if not valid_labels:
            return None

        with plt.rc_context(_MPL_RC):
            fig, ax = plt.subplots(figsize=(5.8, 3.5), dpi=150)
            y_pos = np.arange(len(valid_labels))
            bar_h = 0.35

            ax.barh(y_pos - bar_h / 2, target_vals, bar_h,
                    label=snap.get("ticker", "Target"),
                    color="#6366F1", edgecolor="none")
            ax.barh(y_pos + bar_h / 2, peer_medians, bar_h,
                    label="Peer Median",
                    color="#475569", edgecolor="none")

            ax.set_yticks(y_pos)
            ax.set_yticklabels(valid_labels, fontsize=10)
            ax.legend(loc="lower right", fontsize=9, facecolor="#0B142B",
                      edgecolor="#1E293B", labelcolor="white")
            ax.grid(axis="x", alpha=0.3)
            ax.set_title("Target vs Peer Median", fontsize=13, color="white",
                         pad=12, fontweight="bold")
            fig.tight_layout()

            path = os.path.join(self._charts_dir, "peer_chart.png")
            fig.savefig(path, bbox_inches="tight", facecolor="#0B142B")
            plt.close(fig)
            return path

    def _create_fair_value_chart(self, price, low, mid, high) -> str | None:
        """Horizontal bar showing current price vs fair value range."""
        vals = [price, low, mid, high]
        if any(v is None for v in vals):
            return None
        try:
            price, low, mid, high = [float(v) for v in vals]
        except (TypeError, ValueError):
            return None

        with plt.rc_context(_MPL_RC):
            fig, ax = plt.subplots(figsize=(12.5, 1.9), dpi=150)

            all_vals = [price, low, mid, high]
            x_min = min(all_vals) * 0.82
            x_max = max(all_vals) * 1.18

            # Gradient effect: bear zone + bull zone
            ax.barh(0, mid - low, left=low, height=0.45, color="#EF4444",
                    alpha=0.3, edgecolor="none")
            ax.barh(0, high - mid, left=mid, height=0.45, color="#10B981",
                    alpha=0.3, edgecolor="none")
            # Full range outline
            ax.barh(0, high - low, left=low, height=0.45, color="none",
                    edgecolor="#6366F1", linewidth=2)

            # Mid marker
            ax.plot(mid, 0, marker="D", color="#6366F1", markersize=14, zorder=5)
            ax.annotate(f"Fair Value\n${mid:.2f}", (mid, 0.4),
                        ha="center", va="bottom", fontsize=10, color="#6366F1",
                        fontweight="bold")

            # Current price marker
            ax.plot(price, 0, marker="v", color="#F59E0B", markersize=16, zorder=5)
            ax.annotate(f"Current Price\n${price:.2f}", (price, -0.4),
                        ha="center", va="top", fontsize=10, color="#F59E0B",
                        fontweight="bold")

            # Bear / Bull labels
            ax.annotate(f"Bear: ${low:.2f}", (low, -0.4),
                        ha="center", va="top", fontsize=9, color="#EF4444")
            ax.annotate(f"Bull: ${high:.2f}", (high, -0.4),
                        ha="center", va="top", fontsize=9, color="#10B981")

            ax.set_xlim(x_min, x_max)
            ax.set_ylim(-0.85, 0.85)
            ax.set_yticks([])
            ax.xaxis.set_major_formatter(mticker.FormatStrFormatter("$%.0f"))
            ax.set_title("Current Price vs Fair Value Range", fontsize=12,
                         color="white", pad=10, fontweight="bold")
            ax.grid(axis="x", alpha=0.2)
            fig.tight_layout()

            path = os.path.join(self._charts_dir, "fair_value.png")
            fig.savefig(path, bbox_inches="tight", facecolor="#0B142B")
            plt.close(fig)
            return path

    def _create_sensitivity_heatmap(self, dcf: dict, current_price) -> str | None:
        """5x5 WACC vs terminal growth heatmap."""
        grid = _sg(dcf, "sensitivity_grid", "implied_prices")
        if not grid or not isinstance(grid, dict):
            return None

        try:
            current_price = float(current_price) if current_price else None
        except (TypeError, ValueError):
            current_price = None

        wacc_keys = sorted(grid.keys())
        if not wacc_keys:
            return None

        first_tg_map = grid[wacc_keys[0]]
        if not isinstance(first_tg_map, dict):
            return None
        tg_keys = sorted(first_tg_map.keys())

        data = []
        for wk in wacc_keys:
            row = []
            tg_map = grid.get(wk, {})
            for tk in tg_keys:
                v = tg_map.get(tk)
                row.append(float(v) if v is not None else np.nan)
            data.append(row)

        arr = np.array(data)
        if arr.size == 0:
            return None

        with plt.rc_context(_MPL_RC):
            fig, ax = plt.subplots(figsize=(6.5, 3.3), dpi=150)

            # Color based on distance from current price
            if current_price:
                vmin = current_price * 0.3
                vmax = current_price * 2.0
                cmap = "RdYlGn"
            else:
                vmin, vmax = np.nanmin(arr), np.nanmax(arr)
                cmap = "YlGnBu"

            im = ax.imshow(arr, aspect="auto", cmap=cmap, vmin=vmin, vmax=vmax)

            # Annotate cells
            for i in range(arr.shape[0]):
                for j in range(arr.shape[1]):
                    val = arr[i, j]
                    if np.isnan(val):
                        continue
                    text_color = "black" if val > (vmin + vmax) / 2 else "white"
                    ax.text(j, i, f"${val:.0f}", ha="center", va="center",
                            fontsize=9, color=text_color, fontweight="bold")

            # Format tick labels
            def _fmt_key(k):
                try:
                    return f"{float(k) * 100:.1f}%"
                except (TypeError, ValueError):
                    return str(k)

            ax.set_xticks(range(len(tg_keys)))
            ax.set_xticklabels([_fmt_key(k) for k in tg_keys], fontsize=9)
            ax.set_yticks(range(len(wacc_keys)))
            ax.set_yticklabels([_fmt_key(k) for k in wacc_keys], fontsize=9)

            ax.set_xlabel("Terminal Growth Rate", fontsize=10)
            ax.set_ylabel("WACC", fontsize=10)
            ax.set_title("DCF Sensitivity Analysis", fontsize=12, color="white",
                         pad=10, fontweight="bold")

            cbar = fig.colorbar(im, ax=ax, shrink=0.8, pad=0.02)
            cbar.ax.tick_params(labelsize=9)

            if current_price:
                cbar.ax.axhline(y=current_price, color="#F59E0B", linewidth=2,
                                label="Current")

            fig.tight_layout()

            path = os.path.join(self._charts_dir, "sensitivity.png")
            fig.savefig(path, bbox_inches="tight", facecolor="#0B142B")
            plt.close(fig)
            return path

    def _create_range_chart(self, combined: dict, current_price) -> str | None:
        """Bull/base/bear horizontal bar with labels."""
        low = combined.get("low")
        mid = combined.get("mid")
        high = combined.get("high")

        vals = [low, mid, high, current_price]
        if any(v is None for v in vals):
            return None

        try:
            low, mid, high, cp = [float(v) for v in vals]
        except (TypeError, ValueError):
            return None

        with plt.rc_context(_MPL_RC):
            fig, ax = plt.subplots(figsize=(5.5, 1.4), dpi=150)

            x_min = min(low, cp) * 0.88
            x_max = max(high, cp) * 1.12

            # Gradient bar segments
            ax.barh(0, mid - low, left=low, height=0.5, color="#EF4444",
                    alpha=0.5, edgecolor="none")
            ax.barh(0, high - mid, left=mid, height=0.5, color="#10B981",
                    alpha=0.5, edgecolor="none")

            # Markers
            ax.plot(cp, 0, marker="v", color="#F59E0B", markersize=15, zorder=5)
            ax.plot(mid, 0, marker="D", color="#6366F1", markersize=11, zorder=5)

            ax.annotate(f"Current ${cp:.0f}", (cp, 0.42), ha="center", fontsize=9,
                        color="#F59E0B", fontweight="bold")
            ax.annotate(f"Base ${mid:.0f}", (mid, -0.42), ha="center", va="top",
                        fontsize=9, color="#6366F1", fontweight="bold")

            ax.set_xlim(x_min, x_max)
            ax.set_ylim(-0.65, 0.7)
            ax.set_yticks([])
            ax.xaxis.set_major_formatter(mticker.FormatStrFormatter("$%.0f"))
            ax.set_title("Bull / Base / Bear Range", fontsize=11, color="white",
                         pad=8, fontweight="bold")
            ax.grid(axis="x", alpha=0.2)
            fig.tight_layout()

            path = os.path.join(self._charts_dir, "range.png")
            fig.savefig(path, bbox_inches="tight", facecolor="#0B142B")
            plt.close(fig)
            return path
